from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    """Creates a PowerPoint presentation summarizing the Bhagavad Gita RAG project."""
    prs = Presentation()

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Bhagavad Gita RAG System for Kannada"
    subtitle.text = "A summary of the project"

    # Slide 2: Project Goal
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Project Goal"
    content_shape = slide.shapes.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.text = "To create a Retrieval Augmented Generation (RAG) system for the Bhagavad Gita in Kannada, allowing users to search for verses using natural language queries."
    text_frame.paragraphs[0].font.size = Pt(24)


    # Slide 3: Architecture
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Architecture"
    content_shape = slide.shapes.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()
    p1 = text_frame.add_paragraph()
    p1.text = "Data Loading: Loads Bhagavad Gita verses from a JSON file."
    p1.level = 0
    p1.font.size = Pt(20)
    p2 = text_frame.add_paragraph()
    p2.text = "Embedding: Uses sentence-transformers (paraphrase-multilingual-MiniLM-L12-v2) to create vector embeddings for each verse."
    p2.level = 0
    p2.font.size = Pt(20)
    p3 = text_frame.add_paragraph()
    p3.text = "Retrieval: Encodes the user query and uses cosine similarity to find the most relevant verses."
    p3.level = 0
    p3.font.size = Pt(20)


    # Slide 4: Key Features
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Key Features"
    content_shape = slide.shapes.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()
    features = [
        "Semantic search for Bhagavad Gita verses in Kannada.",
        "Bilingual user interface (Kannada and English).",
        "Text-to-speech functionality to listen to verses and translations.",
        "Adjustable number of search results.",
        "Example queries to guide the user.",
    ]
    for feature in features:
        p = text_frame.add_paragraph()
        p.text = feature
        p.level = 0
        p.font.size = Pt(20)

    # Slide 5: Technology Stack
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Technology Stack"
    content_shape = slide.shapes.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()
    techs = [
        "Streamlit: For the web application UI.",
        "Sentence-Transformers: For semantic search and embeddings.",
        "Scikit-learn: For cosine similarity calculation.",
        "Numpy: For numerical operations.",
        "gTTS: For text-to-speech.",
        "python-pptx: For presentation generation."
    ]
    for tech in techs:
        p = text_frame.add_paragraph()
        p.text = tech
        p.level = 0
        p.font.size = Pt(20)


    # Slide 6: How to Use
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "How to Use"
    content_shape = slide.shapes.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()
    steps = [
        "Run the Streamlit app: streamlit run app.py",
        "Select the language (Kannada or English).",
        "Enter a query in the search box.",
        "Adjust the number of results with the slider.",
        "Click the 'Search' button to see the results.",
    ]
    for step in steps:
        p = text_frame.add_paragraph()
        p.text = step
        p.level = 0
        p.font.size = Pt(20)

    # Slide 7: Conclusion
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Conclusion"
    content_shape = slide.shapes.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.text = "The project successfully implements a functional RAG system for the Bhagavad Gita in Kannada, with a user-friendly web interface and useful features like text-to-speech."
    text_frame.paragraphs[0].font.size = Pt(24)


    prs.save("Bhagavad_Gita_RAG_Summary.pptx")
    print("Presentation 'Bhagavad_Gita_RAG_Summary.pptx' created successfully.")

if __name__ == "__main__":
    create_presentation()
